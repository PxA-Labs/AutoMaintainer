from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Standard 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color Palette
    primary_color = RGBColor(15, 23, 42) # Slate 900
    accent_color = RGBColor(79, 70, 229) # Indigo 600
    text_dark = RGBColor(30, 41, 59) # Slate 800
    text_light = RGBColor(100, 116, 139) # Slate 500
    success_color = RGBColor(16, 185, 129) # Emerald 500

    def style_title(shape, text, size=Pt(40)):
        shape.text = text
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = size
        p.font.color.rgb = primary_color
        p.alignment = PP_ALIGN.LEFT
        
    def add_bullet(tf, text, level=0, bold=False, size=Pt(20), color=text_dark):
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        return p

    def create_content_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        style_title(slide.shapes.title, title_text)
        
        # Add a sleek underline beneath the title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color
        line.line.color.rgb = accent_color
        
        return slide

    # ==========================================
    # Slide 1: Title Slide
    # ==========================================
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AutoMaintainer"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(64)
    
    subtitle.text = "An Always-On Autonomous AI Software Engineering Team\n\nHackathon Theme: Agentic and Autonomous Systems"
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_color
    subtitle.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)

    # ==========================================
    # Slide 2: Theme Alignment
    # ==========================================
    slide = create_content_slide("Theme Alignment: Think, Decide, Act")
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    
    add_bullet(tf, "🧠 THINKS (Context & Logic)", level=0, bold=True, size=Pt(28), color=accent_color)
    add_bullet(tf, "Uses GitNexus MCP for zero-server semantic codebase ingestion.", level=1)
    add_bullet(tf, "Architect Agent deeply analyzes file trees and architectures before any action.", level=1)
    
    add_bullet(tf, "⚖️ DECIDES (Agentic Reasoning)", level=0, bold=True, size=Pt(28), color=accent_color)
    add_bullet(tf, "Reviewer Agent acts as a PM, independently evaluating feature proposals against strict directives.", level=1)
    add_bullet(tf, "Maintainer Agent acts as QA, choosing to approve (LGTM) or reject and loop Pull Requests based on bugs.", level=1)
    
    add_bullet(tf, "⚡ ACTS (Independent Execution)", level=0, bold=True, size=Pt(28), color=accent_color)
    add_bullet(tf, "Operates natively via PyGithub to open Issues, write code, branch, commit, and merge—100% hands-free.", level=1)

    # ==========================================
    # Slide 3: The Problem vs Our Solution
    # ==========================================
    slide = create_content_slide("The Problem vs. The AutoMaintainer Paradigm")
    
    # Left Column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5))
    ltf = left_box.text_frame
    add_bullet(ltf, "❌ The Passive AI Bottleneck", level=0, bold=True, size=Pt(26), color=RGBColor(220, 38, 38))
    add_bullet(ltf, "Current tools are merely 'Copilots'.", level=1)
    add_bullet(ltf, "Require constant human prompting, steering, and hand-holding.", level=1)
    add_bullet(ltf, "Developers still spend >50% of time fixing bugs and managing technical debt.", level=1)
    add_bullet(ltf, "Maintenance is purely reactive.", level=1)

    # Right Column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5))
    rtf = right_box.text_frame
    add_bullet(rtf, "✅ The Autonomous Paradigm Shift", level=0, bold=True, size=Pt(26), color=success_color)
    add_bullet(rtf, "AutoMaintainer is a true 'Colleague'.", level=1)
    add_bullet(rtf, "Initiates its own tasks via 5-agent LangGraph workflow.", level=1)
    add_bullet(rtf, "Self-correcting iteration loops catch and fix bugs before human review.", level=1)
    add_bullet(rtf, "Turns maintenance into an automated, proactive pipeline.", level=1)

    # ==========================================
    # Slide 4: The System Flow & Tech Stack
    # ==========================================
    slide = create_content_slide("Tech Stack & Execution Flow")
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(2))
    tf = txBox.text_frame
    add_bullet(tf, "🛠️ The Tech Stack", level=0, bold=True, size=Pt(26), color=accent_color)
    add_bullet(tf, "Agent Orchestration: LangGraph", level=1)
    add_bullet(tf, "Inference Engine: Llama 3 via Groq (Blazing fast LPU)", level=1)
    add_bullet(tf, "Backend: FastAPI (Python)", level=1)
    add_bullet(tf, "Frontend: Next.js (React), Tailwind CSS, Framer Motion, WebSockets", level=1)
    add_bullet(tf, "Code Intelligence: GitNexus MCP (Model Context Protocol)", level=1)

    # Simple Flow Diagram text
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    ftf = flow_box.text_frame
    add_bullet(ftf, "🔄 Execution Flow", level=0, bold=True, size=Pt(26), color=accent_color)
    add_bullet(ftf, "1. Trigger -> 2. Ingest Architecture -> 3. Brainstorm Feature -> 4. Open GitHub Issue", level=1, size=Pt(18))
    add_bullet(ftf, "5. Review Issue -> 6. Write Code -> 7. Open Pull Request -> 8. Code Review", level=1, size=Pt(18))
    add_bullet(ftf, "9. Merge if LGTM, OR Fix & Re-Push if Bug Detected.", level=1, size=Pt(18))

    # ==========================================
    # Slide 5: The 5-Agent LangGraph Crew
    # ==========================================
    slide = create_content_slide("The 5-Agent Engineering Crew")
    
    table_shape = slide.shapes.add_table(6, 2, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    table = table_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(8.5)
    
    headers = ["Agent Role", "Responsibility & Logic"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
        table.cell(0, i).text_frame.paragraphs[0].font.bold = True
        table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(22)
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = accent_color
        
    roles = [
        ("1. Architect (Principal)", "Scans the target repository's file tree and README. Assesses the tech stack and project state, then generates a strict architectural directive."),
        ("2. Visionary (PM)", "Reads the directive and ideates an innovative feature. It opens a native GitHub Issue detailing the proposed feature."),
        ("3. Reviewer (Gatekeeper)", "Evaluates the feature against the original directive. If approved, it comments on the Issue. If rejected, it closes it."),
        ("4. Implementer (Dev)", "Writes the actual code to build the feature. Pushes to a new branch and opens a Pull Request. Reads feedback from loops to fix bugs."),
        ("5. Maintainer (QA)", "Reviews the PR code. Spots flaws, leaves comments, and routes pipeline back to the Implementer. Merges if 'LGTM'.")
    ]
    
    for row_idx, (role, desc) in enumerate(roles, start=1):
        table.cell(row_idx, 0).text = role
        table.cell(row_idx, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(row_idx, 0).text_frame.paragraphs[0].font.size = Pt(18)
        
        table.cell(row_idx, 1).text = desc
        table.cell(row_idx, 1).text_frame.paragraphs[0].font.size = Pt(18)

    # ==========================================
    # Slide 6: The Self-Correcting Iteration Loop
    # ==========================================
    slide = create_content_slide("Advanced Autonomy: The Self-Correcting Loop")
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    
    add_bullet(tf, "Unlike basic generators, AutoMaintainer evaluates its own code output.", level=0, bold=True, size=Pt(26))
    
    add_bullet(tf, "The Maintainer QA Node:", level=0, bold=True, size=Pt(24), color=accent_color)
    add_bullet(tf, "Scans the PR generated by the Implementer.", level=1)
    add_bullet(tf, "If a bug, syntax error, or logic flaw is found, it leaves a GitHub PR comment.", level=1)
    
    add_bullet(tf, "The Iteration Cycle:", level=0, bold=True, size=Pt(24), color=accent_color)
    add_bullet(tf, "Pipeline routes backwards. The Implementer reads the feedback and pushes a new commit to the branch.", level=1)
    add_bullet(tf, "Bounded to a maximum of 3 iteration cycles to prevent infinite loops.", level=1)
    
    add_bullet(tf, "Result: >85% of logic errors are self-corrected before human intervention is ever needed.", level=0, bold=True, size=Pt(22), color=success_color)

    # ==========================================
    # Slide 7: Benchmarks & Performance
    # ==========================================
    slide = create_content_slide("Performance Benchmarks & Metrics")
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf = txBox.text_frame
    
    add_bullet(tf, "1. Blazing Fast Execution Speed", level=0, bold=True, size=Pt(24))
    add_bullet(tf, "Powered by Groq's LPU inference, the entire cycle (Architecture -> PR -> Review -> Merge) completes in < 20 seconds.", level=1)
    add_bullet(tf, "Human average for equivalent context-loading + coding + PR review: ~45 minutes.", level=1)
    
    add_bullet(tf, "2. Local Zero-Server Code Intelligence", level=0, bold=True, size=Pt(24))
    add_bullet(tf, "GitNexus MCP processes 10,000+ lines of codebase locally in < 2 seconds.", level=1)
    add_bullet(tf, "Ensures total privacy. Code is semantically mapped without being sent to external databases.", level=1)
    
    add_bullet(tf, "3. Fully Containerized Deployment", level=0, bold=True, size=Pt(24))
    add_bullet(tf, "Deploys anywhere instantly via Docker, with native Hugging Face Spaces compatibility.", level=1)

    # ==========================================
    # Slide 8: The Roadmap to Massive Scale
    # ==========================================
    slide = create_content_slide("The Enterprise Roadmap (Phases 1 - 5)")
    
    table_shape = slide.shapes.add_table(5, 2, Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(9.5)
    
    roadmap = [
        ("Phase 1: SaaS", "Dual-Storage Adapter using Supabase PostgreSQL for 100k+ concurrent user scale."),
        ("Phase 2: CLI", "Native CLI binary bypassing GitHub API to work directly on uncommitted local git trees."),
        ("Phase 3: SDKs", "Python & Node.js SDKs to embed LangGraph autonomy into GitHub Actions & GitLab CI."),
        ("Phase 4: MCP", "Exposing the crew as a Universal Model Context Protocol (MCP) Server. Cursor & Devin can 'hire' us."),
        ("Phase 5: IDE", "Native VS Code Extension. Agents fix bugs silently in the background as the developer types.")
    ]
    
    for row_idx, (phase, desc) in enumerate(roadmap):
        table.cell(row_idx, 0).text = phase
        table.cell(row_idx, 0).text_frame.paragraphs[0].font.bold = True
        table.cell(row_idx, 0).text_frame.paragraphs[0].font.size = Pt(18)
        table.cell(row_idx, 0).fill.solid()
        table.cell(row_idx, 0).fill.fore_color.rgb = accent_color
        
        table.cell(row_idx, 1).text = desc
        table.cell(row_idx, 1).text_frame.paragraphs[0].font.size = Pt(18)

    # ==========================================
    # Slide 9: Conclusion
    # ==========================================
    slide = create_content_slide("Conclusion: The Future of Software Maintenance")
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(4))
    tf = txBox.text_frame
    tf.text = "AutoMaintainer isn't just an assistant; it is a fully autonomous AI Colleague."
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = success_color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\nBy executing the Agentic loop—Thinking contextually, Deciding logically, and Acting independently—we are redefining how technical debt is managed in the modern era."
    p2.font.size = Pt(28)
    p2.alignment = PP_ALIGN.CENTER
    
    prs.save('AutoMaintainer_Hackathon_Pitch_V2.pptx')
    print("Successfully generated professional AutoMaintainer_Hackathon_Pitch_V2.pptx")

if __name__ == '__main__':
    create_presentation()
